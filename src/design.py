"""
Design system for the Revolut GP L90D presentation.

Defines colours, typography scale, layout grid, and reusable PPTX helpers.
All slides use these constants for visual consistency.
"""

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─────────────────────────────────────────────────────────────────────────
# COLOUR PALETTE
# ─────────────────────────────────────────────────────────────────────────
# Role-based naming: what a colour MEANS, not what it looks like.

WHITE      = "#FFFFFF"

# Text hierarchy
TEXT_PRIMARY   = "#1A1A2E"   # headings, stat numbers
TEXT_SECONDARY = "#4A5568"   # body copy, card text
TEXT_TERTIARY  = "#94A3B8"   # footnotes, source notes

# Accent: Revolut blue
ACCENT         = "#0B84F6"
ACCENT_BG      = "#EBF5FF"   # light blue card fill

# Semantic
POSITIVE       = "#059669"   # growth, opportunity
POSITIVE_BG    = "#ECFDF5"
NEGATIVE       = "#DC2626"   # decline, risk, cost
NEGATIVE_BG    = "#FEF2F2"
NEUTRAL_BG     = "#F7F8FA"   # informational cards
BORDER         = "#E2E8F0"

# Chart-specific
C_BLUE   = "#0B84F6"
C_BLUE2  = "#60A5FA"
C_BLUE3  = "#93C5FD"
C_RED    = "#DC2626"
C_RED2   = "#F87171"
C_GREEN  = "#059669"


# ─────────────────────────────────────────────────────────────────────────
# TYPOGRAPHY SCALE
# ─────────────────────────────────────────────────────────────────────────
# Minimum body text: 12pt. Footnotes: 9pt (only for source lines).
FONT = "Calibri"

TITLE_SIZE     = 22   # slide title (H1)
SUBTITLE_SIZE  = 13   # slide subtitle
CARD_TITLE     = 12   # insight card heading
BODY_SIZE      = 11   # card body, bullet text
STAT_NUMBER    = 28   # big stat numbers
STAT_LABEL     = 9    # uppercase labels under stats
FOOTNOTE_SIZE  = 8    # source/methodology notes (bottom of slide)
SECTION_LABEL  = 9    # section tags like "GP L90D ANALYSIS"


# ─────────────────────────────────────────────────────────────────────────
# LAYOUT GRID
# ─────────────────────────────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

LM = 0.7              # left margin
RM = 0.7              # right margin
TM = 0.5              # top margin (below brand line)
BM = 0.35             # bottom margin (for footnotes)

CONTENT_W = 13.333 - LM - RM   # ~11.933"
CONTENT_H = 7.5 - TM - BM      # ~6.65"

# Column presets
COL2_W = (CONTENT_W - 0.4) / 2   # two-column width
COL2_GAP = 0.4
COL3_W = (CONTENT_W - 0.6) / 3   # three-column width
COL3_GAP = 0.3

# Footnote y position (consistent across all slides)
FOOTNOTE_Y = 7.05


# ─────────────────────────────────────────────────────────────────────────
# HELPERS
# ─────────────────────────────────────────────────────────────────────────

def I(v):
    """Shorthand for Inches()."""
    return Inches(v)


def rgb(hex_colour):
    """Convert hex string to RGBColor."""
    h = hex_colour.lstrip("#")
    return RGBColor(int(h[:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_text(slide, left, top, width, height, text,
             size=BODY_SIZE, bold=False, colour=TEXT_PRIMARY,
             align=PP_ALIGN.LEFT):
    """Add a simple text box."""
    tb = slide.shapes.add_textbox(I(left), I(top), I(width), I(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(colour)
    p.font.name = FONT
    p.alignment = align
    return tb


def add_rich(slide, left, top, width, height):
    """Add a text box and return its text frame for mixed formatting."""
    tb = slide.shapes.add_textbox(I(left), I(top), I(width), I(height))
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def add_para(tf, text, size=BODY_SIZE, bold=False, colour=TEXT_PRIMARY,
             before=0, after=0, align=PP_ALIGN.LEFT):
    """Append a paragraph to an existing text frame."""
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(colour)
    p.font.name = FONT
    p.space_before = Pt(before)
    p.space_after = Pt(after)
    p.alignment = align
    return p


def add_box(slide, left, top, width, height,
            fill=NEUTRAL_BG, border=None, radius=0.04):
    """Add a rounded rectangle background shape."""
    s = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, I(left), I(top), I(width), I(height))
    s.fill.solid()
    s.fill.fore_color.rgb = rgb(fill)
    if border:
        s.line.color.rgb = rgb(border)
        s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    s.adjustments[0] = radius
    return s


def add_accent_bar(slide, left, top, height, colour=ACCENT, width_pt=4):
    """Vertical accent bar on left edge of a card."""
    s = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, I(left), I(top), Pt(width_pt), I(height))
    s.fill.solid()
    s.fill.fore_color.rgb = rgb(colour)
    s.line.fill.background()
    return s


def add_img(slide, path, left, top, width=None, height=None):
    """Add an image to the slide."""
    kw = {"image_file": str(path), "left": I(left), "top": I(top)}
    if width:
        kw["width"] = I(width)
    if height:
        kw["height"] = I(height)
    return slide.shapes.add_picture(**kw)


def add_brand_line(slide):
    """Add the thin Revolut blue accent line at top of every slide."""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, I(0), I(0), SLIDE_W, Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(ACCENT)
    line.line.fill.background()


def add_page_number(slide, n, total):
    """Add page number at bottom right."""
    add_text(slide, 12.3, 7.1, 0.7, 0.25, f"{n} / {total}",
             size=FOOTNOTE_SIZE, colour=TEXT_TERTIARY, align=PP_ALIGN.RIGHT)


def add_footnote(slide, text):
    """Consistent footnote at the bottom of every content slide."""
    add_text(slide, LM, FOOTNOTE_Y, CONTENT_W, 0.3, text,
             size=FOOTNOTE_SIZE, colour=TEXT_TERTIARY)


def add_title(slide, text):
    """Slide title (H1)."""
    add_text(slide, LM, TM, 12, 0.45, text,
             size=TITLE_SIZE, bold=True, colour=TEXT_PRIMARY)


def add_subtitle(slide, text, y_offset=0.48):
    """Slide subtitle."""
    add_text(slide, LM, TM + y_offset, 12, 0.3, text,
             size=SUBTITLE_SIZE, colour=TEXT_SECONDARY)


def card(slide, left, top, width, height,
         title, body, accent=ACCENT, bg=ACCENT_BG):
    """Insight card with left accent bar, title, and body text."""
    add_box(slide, left, top, width, height, fill=bg, border=BORDER)
    add_accent_bar(slide, left, top, height, colour=accent)
    add_text(slide, left + 0.18, top + 0.1, width - 0.3, 0.25, title,
             size=CARD_TITLE, bold=True, colour=accent)
    add_text(slide, left + 0.18, top + 0.38, width - 0.3, height - 0.48, body,
             size=BODY_SIZE, colour=TEXT_SECONDARY)


def stat_card(slide, left, top, width, value, label,
              detail="", value_colour=TEXT_PRIMARY):
    """Big number stat card with coloured top accent."""
    add_box(slide, left, top, width, 1.0, fill=NEUTRAL_BG)
    # Coloured top edge
    edge = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, I(left), I(top), I(width), Pt(3))
    edge.fill.solid()
    edge.fill.fore_color.rgb = rgb(value_colour)
    edge.line.fill.background()
    # Content
    add_text(slide, left + 0.15, top + 0.1, width - 0.3, 0.4, value,
             size=STAT_NUMBER, bold=True, colour=value_colour)
    add_text(slide, left + 0.15, top + 0.55, width - 0.3, 0.2, label,
             size=STAT_LABEL, bold=True, colour=TEXT_SECONDARY)
    if detail:
        add_text(slide, left + 0.15, top + 0.75, width - 0.3, 0.2, detail,
                 size=FOOTNOTE_SIZE, colour=TEXT_TERTIARY)


def numbered_bullet(slide, num, title, body, y, accent=ACCENT):
    """Numbered bullet with circle + bold title + body text."""
    # Circle
    circ = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, I(LM), I(y), I(0.3), I(0.3))
    circ.fill.solid()
    circ.fill.fore_color.rgb = rgb(accent)
    circ.line.fill.background()
    ctf = circ.text_frame
    ctf.paragraphs[0].text = str(num)
    ctf.paragraphs[0].font.size = Pt(12)
    ctf.paragraphs[0].font.bold = True
    ctf.paragraphs[0].font.color.rgb = rgb(WHITE)
    ctf.paragraphs[0].font.name = FONT
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Mixed-format text
    tf = add_rich(slide, LM + 0.42, y, 11.2, 0.55)
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = title + "  "
    r1.font.size = Pt(BODY_SIZE)
    r1.font.bold = True
    r1.font.color.rgb = rgb(TEXT_PRIMARY)
    r1.font.name = FONT
    r2 = p.add_run()
    r2.text = body
    r2.font.size = Pt(BODY_SIZE)
    r2.font.color.rgb = rgb(TEXT_SECONDARY)
    r2.font.name = FONT


def horizon_column(slide, left, top, width, title, amount,
                   accent, bg, items, body_height=3.4):
    """Roadmap horizon column with header + bullet list."""
    # Header
    hdr = add_box(slide, left, top, width, 0.42, fill=accent)
    add_text(slide, left + 0.12, top + 0.06, width - 0.24, 0.3,
             f"{title}   {amount}",
             size=BODY_SIZE, bold=True, colour=WHITE)
    # Body
    add_box(slide, left, top + 0.47, width, body_height,
            fill=bg, border=BORDER)
    tf = add_rich(slide, left + 0.15, top + 0.58,
                  width - 0.3, body_height - 0.15)
    for item in items:
        add_para(tf, f"\u2022  {item}",
                 size=BODY_SIZE, colour=TEXT_PRIMARY, before=3, after=8)
